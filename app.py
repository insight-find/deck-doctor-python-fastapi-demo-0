from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, ValidationError
from typing import List, Optional
from pptx import Presentation
import io
import json
import re

app = FastAPI(title="PPTX Edit Demo")

class Replacement(BaseModel):
    find: str
    replace: str
    regex: Optional[bool] = False
    ignore_case: Optional[bool] = False

def apply_replacements_to_text(text: str, replacements: List[Replacement]) -> str:
    """Apply the list of replacements to a single text string."""
    if not text:
        return text
    out = text
    for r in replacements:
        if r.regex:
            flags = re.IGNORECASE if r.ignore_case else 0
            try:
                out = re.sub(r.find, r.replace, out, flags=flags)
            except re.error:
                # If the user provided invalid regex, fall back to literal replace
                out = out.replace(r.find, r.replace)
        else:
            if r.ignore_case:
                # do case-insensitive literal replacement using re.escape
                pattern = re.compile(re.escape(r.find), flags=re.IGNORECASE)
                out = pattern.sub(r.replace, out)
            else:
                out = out.replace(r.find, r.replace)
    return out

def replace_text_in_presentation(pres: Presentation, replacements: List[Replacement]):
    """Replace text across shapes, table cells, runs, and notes while preserving run formatting."""
    for slide in pres.slides:
        # shapes (text frames)
        for shape in slide.shapes:
            # tables
            if shape.shape_type == 19 and hasattr(shape, "table"):  # 19 == MSO_SHAPE_TYPE.TABLE
                table = shape.table
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        cell = table.cell(r, c)
                        original = cell.text
                        replaced = apply_replacements_to_text(original, replacements)
                        if replaced != original:
                            # simple assignment is fine for cell.text
                            cell.text = replaced
                continue

            # normal text frames / placeholders
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            # Replace at run-level to preserve formatting
            for para in tf.paragraphs:
                for run in para.runs:
                    original = run.text
                    new_text = apply_replacements_to_text(original, replacements)
                    if new_text != original:
                        run.text = new_text

        # slide notes (if present)
        try:
            notes_slide = slide.notes_slide
            if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                ntf = notes_slide.notes_text_frame
                for para in ntf.paragraphs:
                    for run in para.runs:
                        original = run.text
                        new_text = apply_replacements_to_text(original, replacements)
                        if new_text != original:
                            run.text = new_text
        except Exception:
            # Some slides may not have notes_slide; ignore
            pass

@app.post("/modify-pptx")
async def modify_pptx(
    file: UploadFile = File(...),
    replacements: str = Form(...),  # JSON string representing a list of replacement objects
):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Upload a .pptx file.")

    # parse replacements JSON and validate using Pydantic
    try:
        parsed = json.loads(replacements)
        if not isinstance(parsed, list):
            raise ValueError("replacements JSON must be a list of objects")
        rep_objs = [Replacement.model_validate(item) if not isinstance(item, Replacement) else item for item in parsed]
    except (json.JSONDecodeError, ValidationError, ValueError) as e:
        raise HTTPException(status_code=400, detail=f"Invalid replacements payload: {e}")

    # Read file bytes
    contents = await file.read()
    bio = io.BytesIO(contents)
    try:
        pres = Presentation(bio)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Unable to parse pptx: {e}")

    # perform replacements
    replace_text_in_presentation(pres, rep_objs)

    # Save and return
    out = io.BytesIO()
    pres.save(out)
    out.seek(0)
    headers = {"Content-Disposition": f'attachment; filename="modified-{file.filename}"'}
    return StreamingResponse(
        out,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )
